from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.text import MSO_AUTO_SIZE, PP_PARAGRAPH_ALIGNMENT
from pptx.util import Inches, Pt


def slide1(presentation):
    print("preparing slide 1...")
    lyt = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(lyt)
    title_shape = slide.shapes.title
    subtitle = slide.placeholders[1]
    title_shape.text = "My First PPT Automation"
    subtitle.text = "I DID IT!"
    return slide


def slide2(presentation):
    print("preparing slide 2...")
    lyt = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(lyt)
    shapes = slide.shapes
    title_shapes = shapes.title
    title_shapes.text = "Creating my first slide with tables"
    slide.shapes[1].text_frame.text = 'With colours!'
    left_table_shape = slide.shapes[2]
    left_table_shape.height = Inches(1.25)
    left_table_shape.width = Inches(4.5)
    table = left_table_shape.insert_table(rows=3, cols=2)
    table.height = Inches(1.25)
    table.width = Inches(4.5)
    table.table.columns[0].width = Inches(3.5)
    table.table.columns[1].width = Inches(1.0)
    __table_peach_header(table, 0, 0)
    __table_peach_header(table, 0, 1)
    for row in range(1, 3):
        for col in range(0, 2):
            __table_off_white_cell(table, row, col)
    right_table_shape = slide.shapes[3]
    right_table_shape.height = Inches(2.45)
    right_table_shape.width = Inches(4.5)
    table = right_table_shape.insert_table(rows=6, cols=2)
    table.height = Inches(2.45)
    table.width = Inches(4.5)
    table.table.columns[0].width = Inches(3.5)
    table.table.columns[1].width = Inches(1.0)
    __table_blue_header(table, 0, 0)
    __table_blue_header(table, 0, 1)
    for row in range(1, 6):
        for col in range(0, 2):
            __table_off_white_cell(table, row, col)
    return slide


def slide3(presentation):
    print("preparing slide 3...")
    lyt = presentation.slide_layouts[3]
    slide = presentation.slides.add_slide(lyt)
    shapes = slide.shapes
    title_shapes = shapes.title
    title_shapes.text = "Tables with merged header"
    top_table_shape = slide.shapes[1]
    top_table_shape.height = Inches(1.6)
    top_table_shape.width = Inches(8.49)
    table = top_table_shape.insert_table(rows=4, cols=4)
    table.height = Inches(1.6)
    table.width = Inches(8.50)
    table.table.columns[0].width = Inches(3.25)
    table.table.columns[1].width = Inches(1)
    table.table.columns[2].width = Inches(3.25)
    table.table.columns[3].width = Inches(1)
    for row in range(0, 1):
        for col in range(0, 4):
            __table_peach_header(table, row, col)
    table.table.rows[0].cells[0].merge(table.table.rows[0].cells[3])
    for row in range(1, 2):
        for col in range(0, 4):
            __table_blue_header(table, row, col)
    for row in range(2, 4):
        for col in range(0, 4):
            __table_off_white_cell(table, row, col)
    bottom_table_shape = slide.shapes[2]
    bottom_table_shape.height = Inches(1.6)
    bottom_table_shape.width = Inches(8.49)
    table = bottom_table_shape.insert_table(rows=7, cols=4)
    table.height = Inches(1.6)
    table.width = Inches(8.50)
    table.table.columns[0].width = Inches(3.25)
    table.table.columns[1].width = Inches(1)
    table.table.columns[2].width = Inches(3.25)
    table.table.columns[3].width = Inches(1)
    for row in range(0, 1):
        for col in range(0, 4):
            __table_peach_header(table, row, col)
    table.table.rows[0].cells[0].merge(table.table.rows[0].cells[3])
    for row in range(1, 2):
        for col in range(0, 4):
            __table_blue_header(table, row, col)
    for row in range(2, 7):
        for col in range(0, 4):
            __table_off_white_cell(table, row, col)
    return slide


def slide4(presentation):
    print("preparing slide 4...")
    lyt = presentation.slide_layouts[3]
    slide = presentation.slides.add_slide(lyt)
    shapes = slide.shapes
    title_shapes = shapes.title
    title_shapes.text = "The last slide!"
    top_table_shape = slide.shapes[1]
    top_table_shape.height = Inches(2)
    top_table_shape.width = Inches(8.49)
    table = top_table_shape.insert_table(rows=5, cols=2)
    table.height = Inches(2)
    table.width = Inches(8.49)
    table.table.columns[0].width = Inches(5)
    table.table.columns[1].width = Inches(3.49)
    for row in range(0, 1):
        for col in range(0, 2):
            __table_peach_header(table, row, col)
    table.table.rows[0].cells[0].merge(table.table.rows[0].cells[1])
    for row in range(1, 2):
        for col in range(0, 2):
            __table_blue_header(table, row, col)
    for row in range(2, 5):
        for col in range(0, 2):
            __table_off_white_cell(table, row, col)
    sp = slide.shapes[2]._sp
    sp.getparent().remove(sp)
    return slide


def slide4_pie(presentation, data):
    print("preparing slide 4...")
    lyt = presentation.slide_layouts[3]
    slide = presentation.slides.add_slide(lyt)
    shapes = slide.shapes
    title_shapes = shapes.title
    title_shapes.text = "The last slide!"
    sp = slide.shapes[1]._sp
    sp.getparent().remove(sp)
    sp = slide.shapes[1]._sp
    sp.getparent().remove(sp)

    chart_data = CategoryChartData()
    chart_data.categories = list(data['Type'])
    chart_data.add_series('Series 1', tuple(data['MeasureValue']))
    # add chart to slide --------------------
    x, y, cx, cy = Inches(1), Inches(2), Inches(8.49), Inches(4.09)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.number_format = '0'
    data_labels.position = XL_DATA_LABEL_POSITION.INSIDE_END

    chart.has_title = True
    chart_tt = chart.chart_title
    chart_tt.has_text_frame = True
    chart_ttf = chart_tt.text_frame
    chart_ttf_para = chart_ttf.paragraphs[0]
    chart_ttf_para.text = "Excitement levels"
    chart_ttf_para.font.color.rgb = RGBColor(0x59, 0x59, 0x59)
    chart_ttf_para.font.bold = False
    return slide

def format_header_text(header_text, alignment, table, row_no, cell_no):
    tf = table.table.rows[row_no].cells[cell_no].text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    p = tf.paragraphs[0]
    p.alignment = alignment
    font = p.font
    font.name = "Calibri"
    font.size = Pt(16)
    font.bold = True
    font.color.rgb = RGBColor(0, 0, 0)
    p.text = header_text


def add_data_to_table(data, table, start_row, max_rows, max_cols, start_col=1):
    start_row = start_row - 1
    start_col = start_col - 1
    data_row = 0
    data_col = 0
    for row in range(start_row, start_row + max_rows):
        for col in range(start_col, max_cols):
            tf = table.table.rows[row].cells[col].text_frame
            cell_text = data.iloc[data_row, data_col]
            p = tf.paragraphs[0]
            p.text = str(cell_text)
            if col % 2 == 0:
                p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
            else:
                p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            font = p.font
            font.name = "Calibri"
            font.size = Pt(16)
            font.bold = False
            font.color.rgb = RGBColor(0, 0, 0)
            data_col += 1
        data_row += 1
        data_col = 0
    print("Written data to ppt table...")


def __table_peach_header(table, row_no, cell_no):
    table.table.rows[row_no].cells[cell_no].fill.solid()
    table.table.rows[row_no].cells[cell_no].fill.fore_color.rgb = RGBColor(0xFD, 0xEA, 0xDA)


def __table_blue_header(table, row_no, cell_no):
    table.table.rows[row_no].cells[cell_no].fill.solid()
    table.table.rows[row_no].cells[cell_no].fill.fore_color.rgb = RGBColor(0xD9, 0xF1, 0xFF)


def __table_off_white_cell(table, row_no, cell_no):
    table.table.rows[row_no].cells[cell_no].fill.solid()
    table.table.rows[row_no].cells[cell_no].fill.fore_color.rgb = RGBColor(0xFC, 0xFC, 0xFC)
