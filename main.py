import os

import pandas as pd
from pptx import Presentation
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from q2.pptformatter.PPTFormatter import slide1, slide2, slide3, slide4_pie, format_header_text, add_data_to_table


if __name__ == '__main__':
    print("Starting..")
    prs = Presentation(os.path.join(os.path.abspath(os.path.dirname(__file__)), "../data/template.pptx"))

    slide1(prs)
    print("Done slide 1...")

    slide2 = slide2(prs)
    format_header_text('BASE PEOPLE & CONVERSATIONS', PP_PARAGRAPH_ALIGNMENT.LEFT, slide2.shapes[2], 0, 0)
    format_header_text('#', PP_PARAGRAPH_ALIGNMENT.CENTER, slide2.shapes[2], 0, 1)
    format_header_text('SOURCES', PP_PARAGRAPH_ALIGNMENT.LEFT, slide2.shapes[3], 0, 0)
    format_header_text('%', PP_PARAGRAPH_ALIGNMENT.CENTER, slide2.shapes[3], 0, 1)
    print("Preparing data for slide 2...")
    data = pd.read_csv(os.path.join(os.path.abspath(os.path.dirname(__file__)), "../data/ppt_data.csv"))
    data_base = data.loc[data['Section'] == 'Base']
    data_base = data_base.loc[:, ['Type', 'MeasureValue']].copy().sort_values('MeasureValue')
    add_data_to_table(data_base, slide2.shapes[2], 2, 2, 2)
    data_source = data.loc[data['Section'] == 'Source']
    data_source = data_source.loc[:, ['Type', 'MeasureValue']].copy()
    add_data_to_table(data_source, slide2.shapes[3], 2, 5, 2)
    print("Done slide 2...")

    slide3 = slide3(prs)
    format_header_text('Gender split', PP_PARAGRAPH_ALIGNMENT.CENTER, slide3.shapes[1], 0, 0)
    format_header_text('Before', PP_PARAGRAPH_ALIGNMENT.CENTER, slide3.shapes[1], 1, 0)
    format_header_text('%', PP_PARAGRAPH_ALIGNMENT.CENTER, slide3.shapes[1], 1, 1)
    format_header_text('After', PP_PARAGRAPH_ALIGNMENT.CENTER, slide3.shapes[1], 1, 2)
    format_header_text('%', PP_PARAGRAPH_ALIGNMENT.CENTER, slide3.shapes[1], 1, 3)
    format_header_text('Industry split', PP_PARAGRAPH_ALIGNMENT.CENTER, slide3.shapes[2], 0, 0)
    format_header_text('Before', PP_PARAGRAPH_ALIGNMENT.CENTER, slide3.shapes[2], 1, 0)
    format_header_text('%', PP_PARAGRAPH_ALIGNMENT.CENTER, slide3.shapes[2], 1, 1)
    format_header_text('After', PP_PARAGRAPH_ALIGNMENT.CENTER, slide3.shapes[2], 1, 2)
    format_header_text('%', PP_PARAGRAPH_ALIGNMENT.CENTER, slide3.shapes[2], 1, 3)

    print("Preparing data for slide 3...")
    data_gsb = data.loc[
        (data['Segment'] == 'Before') & (data['Section'] == 'Gender Split') & (data['MeasureType'] == '%')]
    data_gsb = data_gsb.loc[:, ['Type', 'MeasureValue']].copy().reset_index(drop=True).sort_values('MeasureValue',
                                                                                                   ascending=False)
    data_gsa = data.loc[
        (data['Segment'] == 'After') & (data['Section'] == 'Gender Split') & (data['MeasureType'] == '%')]
    data_gsa = data_gsa.loc[:, ['Type', 'MeasureValue']].copy().reset_index(drop=True).sort_values('MeasureValue')
    male_only_before = (data_gsb.loc[data_gsb['Type'] == 'Male Only', 'MeasureValue']).reset_index(drop=True)
    female_only_before = (data_gsb.loc[data_gsb['Type'] == 'Female Only', 'MeasureValue']).reset_index(drop=True)
    male_only_after = (data_gsa.loc[data_gsa['Type'] == 'Male Only', 'MeasureValue']).reset_index(drop=True)
    female_only_after = (data_gsa.loc[data_gsa['Type'] == 'Female Only', 'MeasureValue']).reset_index(drop=True)
    dict1 = {'Before': ['Male Only', 'Female Only'],
             '%B': [item[0] for item in [male_only_before, female_only_before] if len(item) > 0],
             'After': ['Male Only', 'Female Only'],
             '%A': [item[0] for item in [male_only_after, female_only_after] if len(item) > 0]}

    data_gs_final = pd.DataFrame(dict1).reset_index(drop=True)
    add_data_to_table(data_gs_final, slide3.shapes[1], 3, 2, 4)

    data_isb = data.loc[
        (data['Segment'] == 'Before') & (data['Type'] == 'Industry split') & (data['MeasureType'] == '%')]
    data_isb = data_isb.loc[:, ['Content', 'MeasureValue']].copy().sort_values('MeasureValue', ascending=False)
    data_isa = data.loc[
        (data['Segment'] == 'After') & (data['Type'] == 'Industry split') & (data['MeasureType'] == '%')]
    data_isa = data_isa.loc[:, ['Content', 'MeasureValue']].copy().sort_values('MeasureValue', ascending=False)
    add_data_to_table(data_isb, slide3.shapes[2], 3, 5, 2)
    add_data_to_table(data_isa, slide3.shapes[2], 3, 5, 4, 3)
    print("Done slide 3...")

    print("Preparing data for slide 4...")
    data_el = data.loc[(data['Section'] == 'Excitement levels') & (data['MeasureType'] == '%')]
    data_el = data_el.loc[:, ['Type', 'MeasureValue']].copy().sort_values('MeasureValue', ascending=False)
    slide4 = slide4_pie(prs, data_el)
    print("Done slide 4...")
    
    prs.save("output_chart.pptx")  # saving file
    print("Saved PPT...")
